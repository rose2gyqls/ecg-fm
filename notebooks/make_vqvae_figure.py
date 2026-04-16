"""
Beat VQ-VAE figure — VQ-VAE paper style (compact encoder/decoder, codebook top,
straight-through gradient arrow). Single slide PPT.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ── palette (inspired by reference figure) ─────────────────────────────────
C_TEXT      = RGBColor(0x1A, 0x1A, 0x1A)
C_MUTED     = RGBColor(0x55, 0x55, 0x55)
C_TITLE     = RGBColor(0x24, 0x48, 0x9C)
C_BEAT_BG   = RGBColor(0xF3, 0xF6, 0xFC)
C_ENCDEC    = RGBColor(0xE4, 0xEE, 0xFA)
C_ENCDEC_LN = RGBColor(0x55, 0x88, 0xC8)
C_ZE        = RGBColor(0xC9, 0xE2, 0xC2)  # light green cube
C_ZE_DK     = RGBColor(0x8F, 0xC1, 0x88)
C_ZQ        = RGBColor(0xCE, 0xC2, 0xE8)  # light purple cube
C_ZQ_DK     = RGBColor(0x9D, 0x8B, 0xCB)
C_EMB       = RGBColor(0xDD, 0xD2, 0xEE)
C_EMB_LN    = RGBColor(0x7E, 0x69, 0xB4)
C_ARROW     = RGBColor(0x4D, 0x8B, 0xD4)
C_REDARR    = RGBColor(0xD1, 0x43, 0x3A)
C_NODE      = RGBColor(0x4D, 0x8B, 0xD4)
C_QZX       = RGBColor(0xF6, 0xEF, 0xC9)
C_LOSS      = RGBColor(0xFF, 0xF2, 0xCC)


# ── helpers ────────────────────────────────────────────────────────────────
def _set_line_dash(line, val="dash"):
    ln = line._get_or_add_ln()
    prst = etree.SubElement(ln, qn("a:prstDash"))
    prst.set("val", val)


def add_box(slide, x, y, w, h, text, fill, *,
            font_size=10, bold=False, text_color=C_TEXT,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, line_color=None,
            line_width=0.75, dash=False):
    s = slide.shapes.add_shape(shape, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line_color if line_color else RGBColor(0x66, 0x66, 0x66)
    s.line.width = Pt(line_width)
    if dash:
        _set_line_dash(s.line)
    tf = s.text_frame
    tf.margin_left = tf.margin_right = Emu(36000)
    tf.margin_top = tf.margin_bottom = Emu(18000)
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.color.rgb = text_color
        r.font.name = "Calibri"
    return s


def add_text(slide, x, y, w, h, text, *,
             font_size=10, bold=False, italic=False,
             color=C_TEXT, align=PP_ALIGN.CENTER):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run()
        r.text = line
        r.font.size = Pt(font_size)
        r.font.bold = bold
        r.font.italic = italic
        r.font.color.rgb = color
        r.font.name = "Calibri"
    return tb


def add_arrow(slide, x1, y1, x2, y2, *, color=C_ARROW, weight=2.2, dash=False,
              tail="triangle"):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", tail)
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    if dash:
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")
    return conn


def add_curved_arrow(slide, x1, y1, x2, y2, *, color=C_ARROW, weight=2.2,
                     dash=False, tail="triangle"):
    """Elbow / bent connector to suggest a curve."""
    conn = slide.shapes.add_connector(MSO_CONNECTOR.ELBOW, x1, y1, x2, y2)
    conn.line.color.rgb = color
    conn.line.width = Pt(weight)
    ln = conn.line._get_or_add_ln()
    tailEnd = etree.SubElement(ln, qn("a:tailEnd"))
    tailEnd.set("type", tail)
    tailEnd.set("w", "med")
    tailEnd.set("h", "med")
    if dash:
        prst = etree.SubElement(ln, qn("a:prstDash"))
        prst.set("val", "dash")
    return conn


def add_cube(slide, x, y, w, h, depth, fill_front, fill_top, fill_side,
             *, line_color=None):
    """Fake 3D cube with 3 freeform faces."""
    lc = line_color or RGBColor(0x55, 0x55, 0x55)
    # front
    front = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    front.fill.solid(); front.fill.fore_color.rgb = fill_front
    front.line.color.rgb = lc; front.line.width = Pt(0.75)
    # top parallelogram via freeform
    top = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM, x, y - depth, w, depth)
    top.fill.solid(); top.fill.fore_color.rgb = fill_top
    top.line.color.rgb = lc; top.line.width = Pt(0.75)
    # shift parallelogram so slanted side aligns: pptx parallelogram slants right-top.
    # side
    side = slide.shapes.add_shape(MSO_SHAPE.PARALLELOGRAM,
                                  x + w, y - depth + depth / 2,
                                  depth, h)
    # (side visual is approximate — front box is the main readable surface)
    side.fill.solid(); side.fill.fore_color.rgb = fill_side
    side.line.color.rgb = lc; side.line.width = Pt(0.75)
    return front


def add_striped_bar(slide, x, y, w, h, n_stripes, fill, stripe_color,
                    *, line_color=None):
    """Rectangle divided vertically into n_stripes (like codebook entry)."""
    lc = line_color or RGBColor(0x55, 0x55, 0x55)
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    rect.fill.solid(); rect.fill.fore_color.rgb = fill
    rect.line.color.rgb = lc; rect.line.width = Pt(0.75)
    stripe_w = w / n_stripes
    for i in range(1, n_stripes):
        sx = x + int(stripe_w * i)
        line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, sx, y, sx, y + h)
        line.line.color.rgb = stripe_color
        line.line.width = Pt(0.5)
    return rect


# ── build ──────────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])

# Title
add_text(slide, Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.40),
         "Beat VQ-VAE Tokenizer",
         font_size=22, bold=True, color=C_TITLE, align=PP_ALIGN.LEFT)
add_text(slide, Inches(0.3), Inches(0.55), Inches(12.7), Inches(0.30),
         "Per-beat latent → nearest codebook entry (EMA-updated) → reconstruction. "
         "Red arrow = straight-through gradient for encoder.",
         font_size=11, italic=True, color=C_MUTED, align=PP_ALIGN.LEFT)

# ── Codebook (top oval) ───────────────────────────────────────────────────
CB_X = Inches(3.70); CB_Y = Inches(1.05); CB_W = Inches(6.0); CB_H = Inches(1.35)
cb_bg = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, CB_X, CB_Y, CB_W, CB_H)
cb_bg.adjustments[0] = 0.5
cb_bg.fill.solid(); cb_bg.fill.fore_color.rgb = RGBColor(0xFA, 0xF6, 0xFF)
cb_bg.line.color.rgb = C_EMB_LN; cb_bg.line.width = Pt(1.25)

# codebook entries (vertical striped bars)
n_entries_vis = 6
entry_w = Inches(0.55); entry_h = Inches(0.95)
entry_gap = Inches(0.06)
entries_total_w = n_entries_vis * entry_w + (n_entries_vis - 1) * entry_gap
ent_x0 = CB_X + Inches(0.35)
ent_y  = CB_Y + Inches(0.25)

labels = ["e₁", "e₂", "e₃", "e₄", "…", "e_K"]
for i in range(n_entries_vis):
    ex = ent_x0 + i * (entry_w + entry_gap)
    if i == 4:  # ellipsis slot
        add_text(slide, ex, ent_y, entry_w, entry_h, "…",
                 font_size=22, bold=True, color=C_EMB_LN)
        continue
    add_striped_bar(slide, ex, ent_y, entry_w, entry_h,
                    n_stripes=6, fill=C_EMB, stripe_color=C_EMB_LN,
                    line_color=C_EMB_LN)
    add_text(slide, ex, ent_y - Inches(0.28), entry_w, Inches(0.22),
             labels[i], font_size=10, bold=True, color=C_TEXT)

# "Embedding Space" label (right side of oval)
add_text(slide, CB_X + CB_W - Inches(1.35), CB_Y + Inches(0.45),
         Inches(1.30), Inches(0.45),
         "Embedding\nSpace\n(K=512, D=256)",
         font_size=10, bold=True, color=C_EMB_LN)

# ── Main row anchors ───────────────────────────────────────────────────────
MID_Y = Inches(4.35)     # center line for main flow

# ── Input beat ─────────────────────────────────────────────────────────────
IN_X = Inches(0.25); IN_Y = Inches(3.55); IN_W = Inches(1.35); IN_H = Inches(1.60)
add_box(slide, IN_X, IN_Y, IN_W, IN_H, "",
        C_BEAT_BG, shape=MSO_SHAPE.RECTANGLE,
        line_color=RGBColor(0xBB, 0xC7, 0xDD))
# tiny PQRST label inside
add_text(slide, IN_X, IN_Y + Inches(0.20), IN_W, Inches(0.35),
         "P   QRS   T", font_size=11, bold=True, color=C_TITLE)
add_text(slide, IN_X, IN_Y + Inches(0.65), IN_W, Inches(0.40),
         "∿‾|\\_/‾∿", font_size=14, color=C_TEXT)
add_text(slide, IN_X, IN_Y + Inches(1.12), IN_W, Inches(0.40),
         "beat segment\nx ∈ ℝ^(B×1×256)",
         font_size=8, italic=True, color=C_MUTED)
add_text(slide, IN_X, IN_Y + IN_H + Inches(0.05), IN_W, Inches(0.25),
         "Input (per-lead)", font_size=10, bold=True, color=C_TITLE)

# ── Encoder (compact) ──────────────────────────────────────────────────────
EN_X = Inches(1.85); EN_Y = Inches(3.75); EN_W = Inches(1.35); EN_H = Inches(1.25)
add_box(slide, EN_X, EN_Y, EN_W, EN_H,
        "CNN\n(Conv1d ×4)\nBN+GELU",
        C_ENCDEC, font_size=11, bold=True, text_color=C_TITLE,
        line_color=C_ENCDEC_LN, line_width=1.25)
add_text(slide, EN_X, EN_Y + EN_H + Inches(0.08), EN_W, Inches(0.30),
         "Encoder", font_size=11, bold=True, color=C_TITLE)
# brace-like subscript
add_text(slide, EN_X, EN_Y + EN_H + Inches(0.32), EN_W, Inches(0.25),
         "(shared, per-lead)", font_size=9, italic=True, color=C_MUTED)

# ── z_e cube ───────────────────────────────────────────────────────────────
ZE_X = Inches(3.45); ZE_Y = Inches(3.90); ZE_W = Inches(0.85); ZE_H = Inches(1.00)
add_cube(slide, ZE_X, ZE_Y, ZE_W, ZE_H, Inches(0.18),
         fill_front=C_ZE, fill_top=C_ZE_DK, fill_side=C_ZE_DK,
         line_color=RGBColor(0x4F, 0x7A, 0x4A))
add_text(slide, ZE_X - Inches(0.15), ZE_Y + ZE_H + Inches(0.08),
         ZE_W + Inches(0.30), Inches(0.30),
         "z_e(x)  (B, D=256)", font_size=10, bold=True, color=C_TEXT)
# 'D' marker on top of cube
add_text(slide, ZE_X, ZE_Y - Inches(0.45), ZE_W, Inches(0.22),
         "D", font_size=9, italic=True, color=C_MUTED)

# ── node1 (quantization input) ─────────────────────────────────────────────
N1_X = Inches(4.70); N1_Y = Inches(4.26); N1_D = Inches(0.22)
n1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, N1_X, N1_Y, N1_D, N1_D)
n1.fill.solid(); n1.fill.fore_color.rgb = C_NODE
n1.line.color.rgb = C_NODE

# ── q(z|x) — discrete index map ────────────────────────────────────────────
Q_X = Inches(5.30); Q_Y = Inches(4.10); Q_W = Inches(1.20); Q_H = Inches(1.00)
add_box(slide, Q_X, Q_Y, Q_W, Q_H,
        "q(z|x)\n\nindex ∈ {1…K}\nper beat",
        C_QZX, font_size=9, bold=False, text_color=C_TEXT,
        line_color=RGBColor(0xBF, 0xA0, 0x49), line_width=1.0,
        shape=MSO_SHAPE.RECTANGLE)
add_text(slide, Q_X, Q_Y - Inches(0.25), Q_W, Inches(0.22),
         "z  (argmin_k ‖z_e − e_k‖²)",
         font_size=9, italic=True, color=C_MUTED)

# ── node2 (post-lookup) ────────────────────────────────────────────────────
N2_X = Inches(6.75); N2_Y = Inches(4.26)
n2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, N2_X, N2_Y, N1_D, N1_D)
n2.fill.solid(); n2.fill.fore_color.rgb = C_NODE
n2.line.color.rgb = C_NODE

# ── z_q cube ───────────────────────────────────────────────────────────────
ZQ_X = Inches(7.30); ZQ_Y = Inches(3.90); ZQ_W = Inches(0.85); ZQ_H = Inches(1.00)
add_cube(slide, ZQ_X, ZQ_Y, ZQ_W, ZQ_H, Inches(0.18),
         fill_front=C_ZQ, fill_top=C_ZQ_DK, fill_side=C_ZQ_DK,
         line_color=RGBColor(0x6B, 0x59, 0x9E))
add_text(slide, ZQ_X - Inches(0.15), ZQ_Y + ZQ_H + Inches(0.08),
         ZQ_W + Inches(0.30), Inches(0.30),
         "z_q(x)  = e_k*", font_size=10, bold=True, color=C_TEXT)
# small e_k overlay
add_text(slide, ZQ_X, ZQ_Y + Inches(0.20), ZQ_W, Inches(0.30),
         "e_k*", font_size=12, bold=True, color=RGBColor(0x3F, 0x2E, 0x6F))

# ── Decoder ────────────────────────────────────────────────────────────────
DE_X = Inches(8.60); DE_Y = Inches(3.75); DE_W = Inches(1.35); DE_H = Inches(1.25)
add_box(slide, DE_X, DE_Y, DE_W, DE_H,
        "CNN\n(ConvT1d ×4)\nBN+GELU",
        C_ENCDEC, font_size=11, bold=True, text_color=C_TITLE,
        line_color=C_ENCDEC_LN, line_width=1.25)
add_text(slide, DE_X, DE_Y + DE_H + Inches(0.08), DE_W, Inches(0.30),
         "Decoder", font_size=11, bold=True, color=C_TITLE)

# ── Reconstruction ─────────────────────────────────────────────────────────
OUT_X = Inches(10.20); OUT_Y = IN_Y; OUT_W = IN_W; OUT_H = IN_H
add_box(slide, OUT_X, OUT_Y, OUT_W, OUT_H, "",
        C_BEAT_BG, shape=MSO_SHAPE.RECTANGLE,
        line_color=RGBColor(0xBB, 0xC7, 0xDD))
add_text(slide, OUT_X, OUT_Y + Inches(0.20), OUT_W, Inches(0.35),
         "P̂   Q̂R̂Ŝ   T̂", font_size=11, bold=True, color=C_TITLE)
add_text(slide, OUT_X, OUT_Y + Inches(0.65), OUT_W, Inches(0.40),
         "∿‾|\\_/‾∿", font_size=14, color=RGBColor(0x66, 0x66, 0x66))
add_text(slide, OUT_X, OUT_Y + Inches(1.12), OUT_W, Inches(0.40),
         "p(x | z_q)\nx̂ ∈ ℝ^(B×1×256)",
         font_size=8, italic=True, color=C_MUTED)
add_text(slide, OUT_X, OUT_Y + OUT_H + Inches(0.05), OUT_W, Inches(0.25),
         "Reconstruction", font_size=10, bold=True, color=C_TITLE)

# ── Flow arrows (blue) ─────────────────────────────────────────────────────
add_arrow(slide, IN_X + IN_W, MID_Y, EN_X, MID_Y)
add_arrow(slide, EN_X + EN_W, MID_Y, ZE_X, MID_Y)
add_arrow(slide, ZE_X + ZE_W, MID_Y, N1_X, N1_Y + N1_D / 2)
add_arrow(slide, N1_X + N1_D, N1_Y + N1_D / 2, Q_X, MID_Y)
add_arrow(slide, Q_X + Q_W, MID_Y, N2_X, N2_Y + N1_D / 2)
add_arrow(slide, N2_X + N1_D, N2_Y + N1_D / 2, ZQ_X, MID_Y)
add_arrow(slide, ZQ_X + ZQ_W, MID_Y, DE_X, MID_Y)
add_arrow(slide, DE_X + DE_W, MID_Y, OUT_X, MID_Y)

# Codebook → quantization: arrow from bottom of oval to node2
CB_BOTTOM_X = CB_X + CB_W / 2
CB_BOTTOM_Y = CB_Y + CB_H
add_arrow(slide, CB_BOTTOM_X, CB_BOTTOM_Y,
          N2_X + N1_D / 2, N2_Y, color=C_EMB_LN, weight=1.75)
add_text(slide, CB_BOTTOM_X - Inches(1.1), CB_BOTTOM_Y + Inches(0.15),
         Inches(1.20), Inches(0.22),
         "lookup  e_k*", font_size=9, italic=True, color=C_EMB_LN,
         align=PP_ALIGN.RIGHT)

# Also arrow from node1 to codebook (distance computation)
add_arrow(slide, N1_X + N1_D / 2, N1_Y,
          CB_BOTTOM_X - Inches(0.3), CB_BOTTOM_Y,
          color=C_EMB_LN, weight=1.5, dash=True)
add_text(slide, N1_X + Inches(0.10), CB_BOTTOM_Y + Inches(0.15),
         Inches(1.8), Inches(0.22),
         "distance to all e_k", font_size=9, italic=True, color=C_EMB_LN,
         align=PP_ALIGN.LEFT)

# ── Straight-through gradient (red curved arrow z_q → z_e) ─────────────────
# Curve goes below the main flow
GRAD_Y = MID_Y + Inches(0.95)
# z_q bottom → down → left → up to z_e bottom
g1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            ZE_X + ZE_W / 2, GRAD_Y - Inches(0.01),
                            (ZQ_X + ZQ_W / 2) - (ZE_X + ZE_W / 2), Inches(0.02))
g1.fill.solid(); g1.fill.fore_color.rgb = C_REDARR
g1.line.color.rgb = C_REDARR
# verticals
v1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                            ZQ_X + ZQ_W / 2 - Inches(0.01),
                            ZQ_Y + ZQ_H,
                            Inches(0.02), GRAD_Y - (ZQ_Y + ZQ_H))
v1.fill.solid(); v1.fill.fore_color.rgb = C_REDARR
v1.line.color.rgb = C_REDARR
# arrow head into z_e (short upward arrow)
add_arrow(slide, ZE_X + ZE_W / 2, GRAD_Y,
          ZE_X + ZE_W / 2, ZE_Y + ZE_H + Inches(0.02),
          color=C_REDARR, weight=2.2)
# label
add_text(slide, Inches(4.6), GRAD_Y + Inches(0.05), Inches(2.5), Inches(0.28),
         "∇_z L   (straight-through)",
         font_size=11, bold=True, italic=True, color=C_REDARR,
         align=PP_ALIGN.CENTER)

# ── Encoder / Decoder brace labels (already placed under each) ─────────────
# (no extra)

# ── Loss panel (bottom) ────────────────────────────────────────────────────
LY = Inches(6.20); LH = Inches(1.10); LX = Inches(0.25); LW = Inches(12.85)
add_box(slide, LX, LY, LW, LH, "", C_LOSS, dash=True,
        line_color=RGBColor(0xC9, 0xA6, 0x2B))
add_text(slide, LX, LY + Inches(0.06), LW, Inches(0.28),
         "Training Objective    L_total = L_rec + α · L_vq + β · L_fid     "
         "(α = 1.0, β = 0.5)",
         font_size=12, bold=True, color=C_TITLE)

sub_y = LY + Inches(0.40); sub_h = Inches(0.62)
sub_w = Inches(4.10); sub_gap = Inches(0.10); sub_x0 = LX + Inches(0.20)

add_box(slide, sub_x0, sub_y, sub_w, sub_h,
        "L_rec = MSE(x̂, x)   —  waveform reconstruction",
        RGBColor(0xFF, 0xFF, 0xFF), font_size=10,
        shape=MSO_SHAPE.RECTANGLE)

add_box(slide, sub_x0 + sub_w + sub_gap, sub_y, sub_w, sub_h,
        "L_vq = β_c · MSE(sg[z_q], z_e)   —  commitment (EMA-updated codebook, γ=0.99)",
        RGBColor(0xFF, 0xFF, 0xFF), font_size=10,
        shape=MSO_SHAPE.RECTANGLE)

add_box(slide, sub_x0 + 2 * (sub_w + sub_gap), sub_y, sub_w, sub_h,
        "L_fid = MSE(∇x̂, ∇x)   —  preserves P / QRS / T sharpness",
        RGBColor(0xFF, 0xFF, 0xFF), font_size=10,
        shape=MSO_SHAPE.RECTANGLE)

# bottom note
add_text(slide, LX, LY + LH + Inches(0.05), LW, Inches(0.22),
         "Diagnostics: perplexity, codebook usage.    "
         "Inference:  encode(x) → (z_q, indices)  — used as Phase-3 pretrain token IDs.",
         font_size=9, italic=True, color=C_MUTED, align=PP_ALIGN.LEFT)


out = "/home1/irteam/local-node-d/hbkimi/ecg-fm/notebooks/beat_vqvae_figure.pptx"
prs.save(out)
print(f"saved: {out}")
